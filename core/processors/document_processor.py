"""Document Processor Module - IA-Influencer-Agent Platform

Industrial-grade document processing engine for content creators and influencers.
Handles document parsing, analysis, conversion, and AI-powered content extraction.

Author: Fahed Mlaiel <mlaiel@live.de>
Copyright (c) 2025 Fahed Mlaiel. All rights reserved.

⚠️ STRICT COPYRIGHT WARNING - Unauthorized use prohibited ⚠️
This software is proprietary and confidential. Any unauthorized use, copying, 
distribution, or commercialization without explicit written permission is 
strictly prohibited and will result in legal action.
Contact: mlaiel@live.de for licensing inquiries.
================================================================================
"""

import asyncio
import logging
import hashlib
import io
import tempfile
import os
from typing import Dict, Any, List, Optional, Union, BinaryIO, Tuple
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from enum import Enum
import json
import time
import mimetypes

# Document processing imports
try:
    import PyPDF2
    import docx
    import openpyxl
    from pptx import Presentation
    import csv
    import xml.etree.ElementTree as ET
    DOC_LIBS_AVAILABLE = True
except ImportError:
    DOC_LIBS_AVAILABLE = False

# Advanced document processing
try:
    import fitz  # PyMuPDF for better PDF handling
    import python_docx2txt
    import xlrd
    import odfpy
    ADVANCED_DOC_LIBS_AVAILABLE = True
except ImportError:
    ADVANCED_DOC_LIBS_AVAILABLE = False

# OCR and image processing for scanned documents
try:
    import pytesseract
    from PIL import Image
    import cv2
    import numpy as np
    OCR_LIBS_AVAILABLE = True
except ImportError:
    OCR_LIBS_AVAILABLE = False

# Text processing for content analysis
try:
    import nltk
    from nltk.tokenize import sent_tokenize, word_tokenize
    import textstat
    import langdetect
    TEXT_ANALYSIS_AVAILABLE = True
except ImportError:
    TEXT_ANALYSIS_AVAILABLE = False

logger = logging.getLogger(__name__)


class DocumentFormat(str, Enum):
    """
Supported document formats"""

    PDF = "pdf"
    DOCX = "docx"
    DOC = "doc"
    RTF = "rtf"
    TXT = "txt"
    XLSX = "xlsx"
    XLS = "xls"
    CSV = "csv"
    PPTX = "pptx"
    PPT = "ppt"
    ODT = "odt"
    ODS = "ods"
    ODP = "odp"
    HTML = "html"
    XML = "xml"
    JSON = "json"
    MARKDOWN = "markdown"


class DocumentType(str, Enum):
    """Types of documents"""

    TEXT_DOCUMENT = "text_document"
    SPREADSHEET = "spreadsheet"
    PRESENTATION = "presentation"
    PDF_DOCUMENT = "pdf_document"
    WEB_DOCUMENT = "web_document"
    STRUCTURED_DATA = "structured_data"
    UNKNOWN = "unknown"


class ExtractionMode(str, Enum):
    """Document extraction modes"""

    TEXT_ONLY = "text_only"
    STRUCTURED = "structured"
    WITH_FORMATTING = "with_formatting"
    WITH_METADATA = "with_metadata"
    FULL_ANALYSIS = "full_analysis"
    OCR_ENABLED = "ocr_enabled"


class DocumentQuality(str, Enum):
    """Document quality levels"""

    LOW = "low"
    MEDIUM = "medium"
    HIGH = "high"
    EXCELLENT = "excellent"


@dataclass
class DocumentProcessingConfig:
    """Configuration for document processing"""
    extraction_mode: ExtractionMode = ExtractionMode.FULL_ANALYSIS
    enable_ocr: bool = True
    ocr_language: str = "eng"
    enable_text_analysis: bool = True
    enable_structure_analysis: bool = True
    enable_metadata_extraction: bool = True
    max_file_size: int = 100 * 1024 * 1024  # 100MB
    max_pages: int = 1000
    preserve_formatting: bool = True
    extract_images: bool = True
    extract_tables: bool = True
    enable_content_classification: bool = True
    enable_language_detection: bool = True
    enable_summarization: bool = True
    max_summary_length: int = 500
    image_quality_threshold: float = 0.7
    text_confidence_threshold: float = 0.8


@dataclass
class DocumentMetadata:
    """Comprehensive document metadata"""
    filename: Optional[str] = None
    file_size: Optional[int] = None
    format: Optional[DocumentFormat] = None
    document_type: Optional[DocumentType] = None
    mime_type: Optional[str] = None
    created_date: Optional[datetime] = None
    modified_date: Optional[datetime] = None
    author: Optional[str] = None
    title: Optional[str] = None
    subject: Optional[str] = None
    keywords: List[str] = field(default_factory=list)
    language: Optional[str] = None
    page_count: Optional[int] = None
    word_count: Optional[int] = None
    character_count: Optional[int] = None
    encoding: Optional[str] = None
    has_images: bool = False
    has_tables: bool = False
    has_charts: bool = False
    has_hyperlinks: bool = False
    is_password_protected: bool = False
    is_signed: bool = False
    compression_ratio: Optional[float] = None
    security_settings: Dict[str, Any] = field(default_factory=dict)


@dataclass
class DocumentStructure:
    """
Document structure information"""
    sections: List[Dict[str, Any]] = field(default_factory=list)
    headings: List[Dict[str, Any]] = field(default_factory=list)
    paragraphs: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    images: List[Dict[str, Any]] = field(default_factory=list)
    hyperlinks: List[Dict[str, Any]] = field(default_factory=list)
    footnotes: List[Dict[str, Any]] = field(default_factory=list)
    headers: List[Dict[str, Any]] = field(default_factory=list)
    footers: List[Dict[str, Any]] = field(default_factory=list)
    page_layout: Dict[str, Any] = field(default_factory=dict)
    style_information: Dict[str, Any] = field(default_factory=dict)


@dataclass
class ContentAnalysis:
    """
Content analysis results"""
    main_topics: List[str] = field(default_factory=list)
    key_entities: List[Dict[str, Any]] = field(default_factory=list)
    sentiment_score: Optional[float] = None
    sentiment_label: Optional[str] = None
    readability_score: Optional[float] = None
    complexity_level: Optional[str] = None
    content_classification: Optional[str] = None
    summary: Optional[str] = None
    extracted_data: Dict[str, Any] = field(default_factory=dict)
    quality_metrics: Dict[str, float] = field(default_factory=dict)


@dataclass
class DocumentFeatures:
    """
Advanced document features"""
    text_content: Optional[str] = None
    structured_content: Dict[str, Any] = field(default_factory=dict)
    extracted_text: Optional[str] = None
    ocr_text: Optional[str] = None
    content_analysis: Optional[ContentAnalysis] = None
    document_structure: Optional[DocumentStructure] = None
    security_analysis: Dict[str, Any] = field(default_factory=dict)
    accessibility_score: Optional[float] = None
    conversion_quality: Optional[float] = None


@dataclass
class DocumentAnalysisResult:
    """
Result of document analysis"""
    success: bool
    metadata: Optional[DocumentMetadata] = None
    features: Optional[DocumentFeatures] = None
    extracted_content: Optional[str] = None
    structured_data: Optional[Dict[str, Any]] = None
    fingerprint: Optional[str] = None
    content_hash: Optional[str] = None
    tags: List[str] = field(default_factory=list)
    warnings: List[str] = field(default_factory=list)
    errors: List[str] = field(default_factory=list)
    processing_time: float = 0.0
    error_message: Optional[str] = None


class DocumentProcessor:
    """
    📄 ENTERPRISE DOCUMENT PROCESSOR
    
    Industrial-grade document processing engine with advanced content extraction,
    structure analysis, and AI-powered insights for creators and businesses.
    """
    
    def __init__(
        self,
        db_session,
        redis_client,
        config: Optional[DocumentProcessingConfig] = None
    ):
        self.db_session = db_session
        self.redis_client = redis_client
        self.config = config or DocumentProcessingConfig()
        self.logger = logging.getLogger(f"{__name__}.DocumentProcessor")
        
        self._initialized = False
        self._supported_formats = set()
        self._ocr_available = False
        
        if not DOC_LIBS_AVAILABLE:
            self.logger.warning("Core document processing libraries not available")
        
        if not ADVANCED_DOC_LIBS_AVAILABLE:
            self.logger.warning("Advanced document processing libraries not available")
        
        if not OCR_LIBS_AVAILABLE:
            self.logger.warning("OCR libraries not available")
    
    async def initialize(self) -> bool:
        """Initialize the document processor"""
        try:
            # Determine supported formats
            if DOC_LIBS_AVAILABLE:
                self._supported_formats.update([
                    DocumentFormat.PDF, DocumentFormat.DOCX, DocumentFormat.XLSX,
                    DocumentFormat.PPTX, DocumentFormat.CSV, DocumentFormat.TXT
                ])
            
            if ADVANCED_DOC_LIBS_AVAILABLE:
                self._supported_formats.update([
                    DocumentFormat.DOC, DocumentFormat.XLS, DocumentFormat.PPT,
                    DocumentFormat.ODT, DocumentFormat.ODS, DocumentFormat.ODP
                ])
            
            # Check OCR availability
            if OCR_LIBS_AVAILABLE and self.config.enable_ocr:
                try:
                    # Test OCR with a simple image
                    test_image = Image.new('RGB', (100, 30), color='white')
                    pytesseract.image_to_string(test_image)
                    self._ocr_available = True
                    self.logger.info("OCR functionality available")
                except Exception as e:
                    self.logger.warning(f"OCR not available: {e}")
            
            # Initialize text analysis if available
            if TEXT_ANALYSIS_AVAILABLE:
                try:
                    nltk.download('punkt', quiet=True)
                    nltk.download('stopwords', quiet=True)
                except:
                    self.logger.warning("Failed to download NLTK data")
            
            self._initialized = True
            self.logger.info(f"✅ Document processor initialized with {len(self._supported_formats)} supported formats")
            return True
            
        except Exception as e:
            self.logger.error(f"❌ Failed to initialize document processor: {e}")
            return False
    
    async def process(
        self,
        content: Union[str, bytes, BinaryIO, Path],
        options: Optional[Dict[str, Any]] = None,
        metadata: Optional[Dict[str, Any]] = None
    ) -> Dict[str, Any]:
        """
        Process document with comprehensive analysis
        
        Args:
            content: Document content (file path, bytes, or file object)
            options: Processing options
            metadata: Additional metadata
            
        Returns:
            Processing result dictionary
        """
        start_time = time.time()
        options = options or {}
        metadata = metadata or {}
        
        try:
            if not self._initialized:
                await self.initialize()
            
            # Load document
            doc_data, doc_metadata = await self._load_document(content, metadata)
            
            if not doc_data:
                return {
                    "success": False,
                    "error_message": "Failed to load document",
                    "processing_time": time.time() - start_time
                }
            
            # Validate document
            validation_result = await self._validate_document(doc_metadata)
            if not validation_result["valid"]:
                return {
                    "success": False,
                    "error_message": validation_result["reason"],
                    "processing_time": time.time() - start_time
                }
            
            # Extract content based on format
            extracted_content = await self._extract_content(doc_data, doc_metadata)
            
            # Perform structure analysis
            document_structure = None
            if self.config.enable_structure_analysis:
                document_structure = await self._analyze_structure(doc_data, doc_metadata)
            
            # Perform content analysis
            content_analysis = None
            if self.config.enable_text_analysis and extracted_content:
                content_analysis = await self._analyze_content(extracted_content)
            
            # Create document features
            features = DocumentFeatures(
                text_content=extracted_content,
                extracted_text=extracted_content,
                content_analysis=content_analysis,
                document_structure=document_structure
            )
            
            # OCR processing for image-based documents
            if self._ocr_available and self.config.enable_ocr:
                ocr_text = await self._perform_ocr(doc_data, doc_metadata)
                if ocr_text:
                    features.ocr_text = ocr_text
                    if not extracted_content:
                        features.text_content = ocr_text
                        features.extracted_text = ocr_text
            
            # Generate fingerprints
            fingerprint = await self._generate_fingerprint(doc_data)
            content_hash = await self._generate_content_hash(extracted_content or "")
            
            # Generate tags
            tags = await self._generate_tags(doc_metadata, features)
            
            # Structured data extraction
            structured_data = await self._extract_structured_data(doc_data, doc_metadata, features)
            
            # Create analysis result
            analysis_result = DocumentAnalysisResult(
                success=True,
                metadata=doc_metadata,
                features=features,
                extracted_content=extracted_content,
                structured_data=structured_data,
                fingerprint=fingerprint,
                content_hash=content_hash,
                tags=tags,
                processing_time=time.time() - start_time
            )
            
            return {
                "success": True,
                "extracted_content": extracted_content,
                "structured_data": structured_data,
                "analysis_result": analysis_result.__dict__,
                "metadata": doc_metadata.__dict__,
                "quality_metrics": {
                    "accessibility_score": features.accessibility_score,
                    "conversion_quality": features.conversion_quality,
                    "content_quality": content_analysis.quality_metrics if content_analysis else None
                },
                "tags": tags,
                "processing_time": time.time() - start_time
            }
            
        except Exception as e:
            self.logger.error(f"Document processing failed: {str(e)}")
            return {
                "success": False,
                "error_message": str(e),
                "processing_time": time.time() - start_time
            }
    
    async def _load_document(
        self,
        content: Union[str, bytes, BinaryIO, Path],
        metadata: Dict[str, Any]
    ) -> Tuple[Optional[bytes], Optional[DocumentMetadata]]:
        """Load document data and extract basic metadata"""
        try:
            doc_data = None
            filename = None
            file_size = None
            
            # Handle different input types
            if isinstance(content, (str, Path)):
                # File path
                file_path = Path(content)
                if not file_path.exists():
                    raise FileNotFoundError(f"Document file not found: {file_path}")
                
                filename = file_path.name
                file_size = file_path.stat().st_size
                
                # Check file size
                if file_size > self.config.max_file_size:
                    raise ValueError(f"File size ({file_size}) exceeds maximum ({self.config.max_file_size})")
                
                with open(file_path, 'rb') as f:
                    doc_data = f.read()
                    
            elif isinstance(content, bytes):
                doc_data = content
                file_size = len(doc_data)
                filename = metadata.get('filename', 'unknown_document')
                
            else:
                # File object
                doc_data = content.read()
                file_size = len(doc_data)
                filename = getattr(content, 'name', metadata.get('filename', 'unknown_document'))
            
            # Determine format and mime type
            doc_format, mime_type = await self._detect_format(doc_data, filename)
            
            # Extract basic metadata
            doc_metadata = DocumentMetadata(
                filename=filename,
                file_size=file_size,
                format=doc_format,
                mime_type=mime_type,
                created_date=datetime.now()
            )
            
            # Determine document type
            doc_metadata.document_type = await self._determine_document_type(doc_format)
            
            return doc_data, doc_metadata
            
        except Exception as e:
            self.logger.error(f"Failed to load document: {e}")
            return None, None
    
    async def _detect_format(self, data: bytes, filename: str) -> Tuple[Optional[DocumentFormat], Optional[str]]:
        """Detect document format and MIME type"""
        try:
            # First try MIME type detection
            mime_type, _ = mimetypes.guess_type(filename)
            
            # Check file signature (magic bytes)
            if data[:4] == b'%PDF':
                return DocumentFormat.PDF, 'application/pdf'
            elif data[:2] == b'PK':  # ZIP-based formats
                if b'word/' in data[:1000] or filename.lower().endswith('.docx'):
                    return DocumentFormat.DOCX, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                elif b'xl/' in data[:1000] or filename.lower().endswith('.xlsx'):
                    return DocumentFormat.XLSX, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
                elif b'ppt/' in data[:1000] or filename.lower().endswith('.pptx'):
                    return DocumentFormat.PPTX, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            elif data[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':  # OLE format
                if filename.lower().endswith('.doc'):
                    return DocumentFormat.DOC, 'application/msword'
                elif filename.lower().endswith('.xls'):
                    return DocumentFormat.XLS, 'application/vnd.ms-excel'
                elif filename.lower().endswith('.ppt'):
                    return DocumentFormat.PPT, 'application/vnd.ms-powerpoint'
            
            # Fallback to filename extension
            extension = Path(filename).suffix.lower()
            format_map = {
                '.pdf': (DocumentFormat.PDF, 'application/pdf'),
                '.docx': (DocumentFormat.DOCX, 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'),
                '.doc': (DocumentFormat.DOC, 'application/msword'),
                '.xlsx': (DocumentFormat.XLSX, 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'),
                '.xls': (DocumentFormat.XLS, 'application/vnd.ms-excel'),
                '.pptx': (DocumentFormat.PPTX, 'application/vnd.openxmlformats-officedocument.presentationml.presentation'),
                '.ppt': (DocumentFormat.PPT, 'application/vnd.ms-powerpoint'),
                '.txt': (DocumentFormat.TXT, 'text/plain'),
                '.csv': (DocumentFormat.CSV, 'text/csv'),
                '.rtf': (DocumentFormat.RTF, 'application/rtf'),
                '.odt': (DocumentFormat.ODT, 'application/vnd.oasis.opendocument.text'),
                '.ods': (DocumentFormat.ODS, 'application/vnd.oasis.opendocument.spreadsheet'),
                '.odp': (DocumentFormat.ODP, 'application/vnd.oasis.opendocument.presentation'),
                '.html': (DocumentFormat.HTML, 'text/html'),
                '.xml': (DocumentFormat.XML, 'application/xml'),
                '.json': (DocumentFormat.JSON, 'application/json'),
                '.md': (DocumentFormat.MARKDOWN, 'text/markdown')
            }
            
            if extension in format_map:
                return format_map[extension]
            
            return None, mime_type
            
        except Exception as e:
            self.logger.error(f"Format detection failed: {e}")
            return None, None
    
    async def _determine_document_type(self, doc_format: Optional[DocumentFormat]) -> DocumentType:
        """Determine document type from format"""
        if not doc_format:
            return DocumentType.UNKNOWN
        
        type_map = {
            DocumentFormat.PDF: DocumentType.PDF_DOCUMENT,
            DocumentFormat.DOCX: DocumentType.TEXT_DOCUMENT,
            DocumentFormat.DOC: DocumentType.TEXT_DOCUMENT,
            DocumentFormat.RTF: DocumentType.TEXT_DOCUMENT,
            DocumentFormat.TXT: DocumentType.TEXT_DOCUMENT,
            DocumentFormat.ODT: DocumentType.TEXT_DOCUMENT,
            DocumentFormat.XLSX: DocumentType.SPREADSHEET,
            DocumentFormat.XLS: DocumentType.SPREADSHEET,
            DocumentFormat.CSV: DocumentType.SPREADSHEET,
            DocumentFormat.ODS: DocumentType.SPREADSHEET,
            DocumentFormat.PPTX: DocumentType.PRESENTATION,
            DocumentFormat.PPT: DocumentType.PRESENTATION,
            DocumentFormat.ODP: DocumentType.PRESENTATION,
            DocumentFormat.HTML: DocumentType.WEB_DOCUMENT,
            DocumentFormat.XML: DocumentType.STRUCTURED_DATA,
            DocumentFormat.JSON: DocumentType.STRUCTURED_DATA,
            DocumentFormat.MARKDOWN: DocumentType.TEXT_DOCUMENT
        }
        
        return type_map.get(doc_format, DocumentType.UNKNOWN)
    
    async def _validate_document(self, metadata: DocumentMetadata) -> Dict[str, Any]:
        """
Validate document against configuration constraints"""
        if not metadata.format or metadata.format not in self._supported_formats:
            return {
                "valid": False,
                "reason": f"Unsupported document format: {metadata.format}"
            }
        
        if metadata.file_size and metadata.file_size > self.config.max_file_size:
            return {
                "valid": False,
                "reason": f"File size ({metadata.file_size}) exceeds maximum ({self.config.max_file_size})"
            }
        
        return {"valid": True}
    
    async def _extract_content(self, data: bytes, metadata: DocumentMetadata) -> Optional[str]:
        """Extract text content from document based on format"""
        try:
            if metadata.format == DocumentFormat.PDF:
                return await self._extract_pdf_content(data)
            elif metadata.format == DocumentFormat.DOCX:
                return await self._extract_docx_content(data)
            elif metadata.format == DocumentFormat.DOC:
                return await self._extract_doc_content(data)
            elif metadata.format == DocumentFormat.XLSX:
                return await self._extract_xlsx_content(data)
            elif metadata.format == DocumentFormat.XLS:
                return await self._extract_xls_content(data)
            elif metadata.format == DocumentFormat.PPTX:
                return await self._extract_pptx_content(data)
            elif metadata.format == DocumentFormat.CSV:
                return await self._extract_csv_content(data)
            elif metadata.format == DocumentFormat.TXT:
                return await self._extract_txt_content(data)
            elif metadata.format == DocumentFormat.RTF:
                return await self._extract_rtf_content(data)
            elif metadata.format == DocumentFormat.HTML:
                return await self._extract_html_content(data)
            elif metadata.format == DocumentFormat.XML:
                return await self._extract_xml_content(data)
            elif metadata.format == DocumentFormat.JSON:
                return await self._extract_json_content(data)
            elif metadata.format == DocumentFormat.MARKDOWN:
                return await self._extract_markdown_content(data)
            else:
                self.logger.warning(f"No extraction method for format: {metadata.format}")
                return None
                
        except Exception as e:
            self.logger.error(f"Content extraction failed: {e}")
            return None
    
    async def _extract_pdf_content(self, data: bytes) -> Optional[str]:
        """Extract text from PDF document"""
        try:
            # Try PyMuPDF first (better quality)
            if ADVANCED_DOC_LIBS_AVAILABLE:
                try:
                    doc = fitz.open(stream=data, filetype="pdf")
                    text_content = []
                    
                    for page_num in range(min(len(doc), self.config.max_pages)):
                        page = doc[page_num]
                        text = page.get_text()
                        if text.strip():
                            text_content.append(text)
                    
                    doc.close()
                    return '\n\n'.join(text_content)
                except Exception as e:
                    self.logger.warning(f"PyMuPDF extraction failed: {e}")
            
            # Fallback to PyPDF2
            if DOC_LIBS_AVAILABLE:
                pdf_reader = PyPDF2.PdfReader(io.BytesIO(data))
                text_content = []
                
                for page_num in range(min(len(pdf_reader.pages), self.config.max_pages)):
                    page = pdf_reader.pages[page_num]
                    text = page.extract_text()
                    if text.strip():
                        text_content.append(text)
                
                return '\n\n'.join(text_content)
            
            return None
            
        except Exception as e:
            self.logger.error(f"PDF extraction failed: {e}")
            return None
    
    async def _extract_docx_content(self, data: bytes) -> Optional[str]:
        """Extract text from DOCX document"""
        try:
            if not DOC_LIBS_AVAILABLE:
                return None
            
            doc = docx.Document(io.BytesIO(data))
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            # Extract table content
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_content.append(' | '.join(row_text))
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            self.logger.error(f"DOCX extraction failed: {e}")
            return None
    
    async def _extract_doc_content(self, data: bytes) -> Optional[str]:
        """Extract text from DOC document"""
        try:
            if ADVANCED_DOC_LIBS_AVAILABLE:
                # Save to temporary file for processing
                with tempfile.NamedTemporaryFile(suffix='.doc', delete=False) as tmp_file:
                    tmp_file.write(data)
                    tmp_file.flush()
                    
                    try:
                        text = python_docx2txt.process(tmp_file.name)
                        return text
                    finally:
                        os.unlink(tmp_file.name)
            
            return None
            
        except Exception as e:
            self.logger.error(f"DOC extraction failed: {e}")
            return None
    
    async def _extract_xlsx_content(self, data: bytes) -> Optional[str]:
        """Extract text from XLSX spreadsheet"""
        try:
            if not DOC_LIBS_AVAILABLE:
                return None
            
            workbook = openpyxl.load_workbook(io.BytesIO(data), data_only=True)
            text_content = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_content.append(f"=== {sheet_name} ===")
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = []
                    for cell_value in row:
                        if cell_value is not None:
                            row_text.append(str(cell_value))
                    if row_text:
                        text_content.append(' | '.join(row_text))
            
            return '\n'.join(text_content)
            
        except Exception as e:
            self.logger.error(f"XLSX extraction failed: {e}")
            return None
    
    async def _extract_xls_content(self, data: bytes) -> Optional[str]:
        """Extract text from XLS spreadsheet"""
        try:
            if ADVANCED_DOC_LIBS_AVAILABLE:
                # Save to temporary file
                with tempfile.NamedTemporaryFile(suffix='.xls', delete=False) as tmp_file:
                    tmp_file.write(data)
                    tmp_file.flush()
                    
                    try:
                        workbook = xlrd.open_workbook(tmp_file.name)
                        text_content = []
                        
                        for sheet_index in range(workbook.nsheets):
                            sheet = workbook.sheet_by_index(sheet_index)
                            text_content.append(f"=== {sheet.name} ===")
                            
                            for row_index in range(sheet.nrows):
                                row_text = []
                                for col_index in range(sheet.ncols):
                                    cell_value = sheet.cell_value(row_index, col_index)
                                    if cell_value:
                                        row_text.append(str(cell_value))
                                if row_text:
                                    text_content.append(' | '.join(row_text))
                        
                        return '\n'.join(text_content)
                    finally:
                        os.unlink(tmp_file.name)
            
            return None
            
        except Exception as e:
            self.logger.error(f"XLS extraction failed: {e}")
            return None
    
    async def _extract_pptx_content(self, data: bytes) -> Optional[str]:
        """Extract text from PPTX presentation"""
        try:
            if not DOC_LIBS_AVAILABLE:
                return None
            
            presentation = Presentation(io.BytesIO(data))
            text_content = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                text_content.append(f"=== Slide {slide_num} ===")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text)
            
            return '\n\n'.join(text_content)
            
        except Exception as e:
            self.logger.error(f"PPTX extraction failed: {e}")
            return None
    
    async def _extract_csv_content(self, data: bytes) -> Optional[str]:
        """Extract text from CSV file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    text = data.decode(encoding)
                    break
                except UnicodeDecodeError:
                    continue
            else:
                return None
            
            # Parse CSV
            csv_reader = csv.reader(io.StringIO(text))
            text_content = []
            
            for row in csv_reader:
                if row:
                    text_content.append(' | '.join(str(cell) for cell in row))
            
            return '\n'.join(text_content)
            
        except Exception as e:
            self.logger.error(f"CSV extraction failed: {e}")
            return None
    
    async def _extract_txt_content(self, data: bytes) -> Optional[str]:
        """Extract text from plain text file"""
        try:
            # Try different encodings
            for encoding in ['utf-8', 'latin-1', 'cp1252']:
                try:
                    return data.decode(encoding)
                except UnicodeDecodeError:
                    continue
            
            # Fallback with error handling
            return data.decode('utf-8', errors='ignore')
            
        except Exception as e:
            self.logger.error(f"TXT extraction failed: {e}")
            return None
    
    async def _extract_rtf_content(self, data: bytes) -> Optional[str]:
        """Extract text from RTF document"""
        try:
            # Basic RTF parsing (simplified)
            text = data.decode('utf-8', errors='ignore')
            
            # Remove RTF control codes (basic)
            import re
            # Remove control words
            text = re.sub(r'\\[a-z]+\d*', '', text)
            # Remove special characters
            text = re.sub(r'[{}\\]', '', text)
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text)
            
            return text.strip()
            
        except Exception as e:
            self.logger.error(f"RTF extraction failed: {e}")
            return None
    
    async def _extract_html_content(self, data: bytes) -> Optional[str]:
        """Extract text from HTML document"""
        try:
            text = data.decode('utf-8', errors='ignore')
            
            # Basic HTML tag removal
            import re
            # Remove script and style elements
            text = re.sub(r'<(script|style)[^>]*>.*?</\1>', '', text, flags=re.DOTALL)
            # Remove HTML tags
            text = re.sub(r'<[^>]+>', '', text)
            # Decode HTML entities
            import html
            text = html.unescape(text)
            # Clean up whitespace
            text = re.sub(r'\s+', ' ', text)
            
            return text.strip()
            
        except Exception as e:
            self.logger.error(f"HTML extraction failed: {e}")
            return None
    
    async def _extract_xml_content(self, data: bytes) -> Optional[str]:
        """Extract text from XML document"""
        try:
            text = data.decode('utf-8', errors='ignore')
            root = ET.fromstring(text)
            
            # Extract all text content
            text_content = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    text_content.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    text_content.append(elem.tail.strip())
            
            return '\n'.join(text_content)
            
        except Exception as e:
            self.logger.error(f"XML extraction failed: {e}")
            return None
    
    async def _extract_json_content(self, data: bytes) -> Optional[str]:
        """Extract text from JSON document"""
        try:
            text = data.decode('utf-8', errors='ignore')
            json_data = json.loads(text)
            
            # Extract text values recursively
            def extract_text_values(obj):
                text_values = []
                if isinstance(obj, dict):
                    for value in obj.values():
                        text_values.extend(extract_text_values(value))
                elif isinstance(obj, list):
                    for item in obj:
                        text_values.extend(extract_text_values(item))
                elif isinstance(obj, str):
                    text_values.append(obj)
                else:
                    text_values.append(str(obj))
                return text_values
            
            text_values = extract_text_values(json_data)
            return '\n'.join(text_values)
            
        except Exception as e:
            self.logger.error(f"JSON extraction failed: {e}")
            return None
    
    async def _extract_markdown_content(self, data: bytes) -> Optional[str]:
        """Extract text from Markdown document"""
        try:
            text = data.decode('utf-8', errors='ignore')
            
            # Basic Markdown parsing (remove formatting)
            import re
            # Remove headers
            text = re.sub(r'^#+\s*', '', text, flags=re.MULTILINE)
            # Remove emphasis
            text = re.sub(r'[*_]{1,2}([^*_]+)[*_]{1,2}', r'\1', text)
            # Remove links
            text = re.sub(r'\[([^\]]+)\]\([^)]+\)', r'\1', text)
            # Remove code blocks
            text = re.sub(r'```[^`]*```', '', text, flags=re.DOTALL)
            text = re.sub(r'`([^`]+)`', r'\1', text)
            # Remove horizontal rules
            text = re.sub(r'^-{3,}$', '', text, flags=re.MULTILINE)
            
            return text.strip()
            
        except Exception as e:
            self.logger.error(f"Markdown extraction failed: {e}")
            return None
    
    async def _analyze_structure(self, data: bytes, metadata: DocumentMetadata) -> Optional[DocumentStructure]:
        """Analyze document structure"""
        try:
            structure = DocumentStructure()
            
            if metadata.format == DocumentFormat.PDF:
                structure = await self._analyze_pdf_structure(data)
            elif metadata.format == DocumentFormat.DOCX:
                structure = await self._analyze_docx_structure(data)
            elif metadata.format == DocumentFormat.XLSX:
                structure = await self._analyze_xlsx_structure(data)
            elif metadata.format == DocumentFormat.PPTX:
                structure = await self._analyze_pptx_structure(data)
            
            return structure
            
        except Exception as e:
            self.logger.error(f"Structure analysis failed: {e}")
            return DocumentStructure()
    
    async def _analyze_pdf_structure(self, data: bytes) -> DocumentStructure:
        """Analyze PDF document structure"""
        structure = DocumentStructure()
        
        try:
            if ADVANCED_DOC_LIBS_AVAILABLE:
                doc = fitz.open(stream=data, filetype="pdf")
                
                for page_num in range(len(doc)):
                    page = doc[page_num]
                    
                    # Get page layout
                    structure.page_layout[f"page_{page_num + 1}"] = {
                        "width": page.rect.width,
                        "height": page.rect.height,
                        "rotation": page.rotation
                    }
                    
                    # Extract text blocks
                    blocks = page.get_text("dict")
                    for block in blocks.get("blocks", []):
                        if "lines" in block:
                            for line in block["lines"]:
                                for span in line["spans"]:
                                    if span["text"].strip():
                                        structure.paragraphs.append({
                                            "text": span["text"],
                                            "page": page_num + 1,
                                            "bbox": span["bbox"],
                                            "font": span["font"],
                                            "size": span["size"]
                                        })
                    
                    # Extract images
                    images = page.get_images()
                    for img_index, img in enumerate(images):
                        structure.images.append({
                            "page": page_num + 1,
                            "index": img_index,
                            "xref": img[0]
                        })
                
                doc.close()
        
        except Exception as e:
            self.logger.error(f"PDF structure analysis failed: {e}")
        
        return structure
    
    async def _analyze_docx_structure(self, data: bytes) -> DocumentStructure:
        """Analyze DOCX document structure"""
        structure = DocumentStructure()
        
        try:
            if DOC_LIBS_AVAILABLE:
                doc = docx.Document(io.BytesIO(data))
                
                # Extract headings and paragraphs
                for para in doc.paragraphs:
                    if para.text.strip():
                        para_info = {
                            "text": para.text,
                            "style": para.style.name if para.style else None
                        }
                        
                        if para.style and 'Heading' in para.style.name:
                            structure.headings.append(para_info)
                        else:
                            structure.paragraphs.append(para_info)
                
                # Extract tables
                for table_index, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    
                    structure.tables.append({
                        "index": table_index,
                        "data": table_data,
                        "rows": len(table.rows),
                        "cols": len(table.columns) if table.rows else 0
                    })
        
        except Exception as e:
            self.logger.error(f"DOCX structure analysis failed: {e}")
        
        return structure
    
    async def _analyze_xlsx_structure(self, data: bytes) -> DocumentStructure:
        """Analyze XLSX spreadsheet structure"""
        structure = DocumentStructure()
        
        try:
            if DOC_LIBS_AVAILABLE:
                workbook = openpyxl.load_workbook(io.BytesIO(data))
                
                for sheet_name in workbook.sheetnames:
                    sheet = workbook[sheet_name]
                    
                    # Get sheet dimensions
                    max_row = sheet.max_row
                    max_col = sheet.max_column
                    
                    sheet_info = {
                        "name": sheet_name,
                        "max_row": max_row,
                        "max_col": max_col,
                        "used_range": f"A1:{sheet.cell(max_row, max_col).coordinate if max_row and max_col else 'A1'}"
                    }
                    
                    structure.sections.append(sheet_info)
        
        except Exception as e:
            self.logger.error(f"XLSX structure analysis failed: {e}")
        
        return structure
    
    async def _analyze_pptx_structure(self, data: bytes) -> DocumentStructure:
        """Analyze PPTX presentation structure"""
        structure = DocumentStructure()
        
        try:
            if DOC_LIBS_AVAILABLE:
                presentation = Presentation(io.BytesIO(data))
                
                for slide_index, slide in enumerate(presentation.slides):
                    slide_info = {
                        "index": slide_index + 1,
                        "shapes": len(slide.shapes),
                        "layout": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else "Unknown"
                    }
                    
                    structure.sections.append(slide_info)
                    
                    # Extract text shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            structure.paragraphs.append({
                                "text": shape.text,
                                "slide": slide_index + 1,
                                "shape_type": str(shape.shape_type)
                            })
        
        except Exception as e:
            self.logger.error(f"PPTX structure analysis failed: {e}")
        
        return structure
    
    async def _analyze_content(self, text: str) -> Optional[ContentAnalysis]:
        """Analyze document content"""
        try:
            analysis = ContentAnalysis()
            
            if not TEXT_ANALYSIS_AVAILABLE:
                return analysis
            
            # Basic text statistics
            word_count = len(text.split())
            
            # Language detection
            try:
                language = langdetect.detect(text)
            except:
                language = "unknown"
            
            # Extract topics (simplified)
            analysis.main_topics = await self._extract_topics(text)
            
            # Sentiment analysis (basic)
            analysis.sentiment_score, analysis.sentiment_label = await self._analyze_sentiment(text)
            
            # Readability analysis
            try:
                analysis.readability_score = textstat.flesch_reading_ease(text)
                
                if analysis.readability_score >= 90:
                    analysis.complexity_level = "very_easy"
                elif analysis.readability_score >= 80:
                    analysis.complexity_level = "easy"
                elif analysis.readability_score >= 70:
                    analysis.complexity_level = "fairly_easy"
                elif analysis.readability_score >= 60:
                    analysis.complexity_level = "standard"
                elif analysis.readability_score >= 50:
                    analysis.complexity_level = "fairly_difficult"
                elif analysis.readability_score >= 30:
                    analysis.complexity_level = "difficult"
                else:
                    analysis.complexity_level = "very_difficult"
            except:
                analysis.readability_score = 50.0
                analysis.complexity_level = "standard"
            
            # Content classification
            analysis.content_classification = await self._classify_content(text)
            
            # Generate summary
            if self.config.enable_summarization and word_count > 100:
                analysis.summary = await self._generate_summary(text)
            
            # Quality metrics
            analysis.quality_metrics = await self._calculate_quality_metrics(text)
            
            return analysis
            
        except Exception as e:
            self.logger.error(f"Content analysis failed: {e}")
            return ContentAnalysis()
    
    async def _extract_topics(self, text: str) -> List[str]:
        """Extract main topics from text"""
        try:
            # Simplified topic extraction using keyword frequency
            if not TEXT_ANALYSIS_AVAILABLE:
                return []
            
            words = word_tokenize(text.lower())
            
            # Remove stopwords
            try:
                from nltk.corpus import stopwords
                stop_words = set(stopwords.words('english'))
            except:
                stop_words = set(['the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'by'])
            
            # Filter words
            filtered_words = [word for word in words if word.isalpha() and len(word) > 3 and word not in stop_words]
            
            # Get most frequent words as topics
            from collections import Counter
            word_freq = Counter(filtered_words)
            topics = [word for word, freq in word_freq.most_common(5)]
            
            return topics
            
        except Exception as e:
            self.logger.error(f"Topic extraction failed: {e}")
            return []
    
    async def _analyze_sentiment(self, text: str) -> Tuple[float, str]:
        """Analyze sentiment of text"""
        try:
            # Basic sentiment analysis
            positive_words = ['good', 'great', 'excellent', 'amazing', 'wonderful', 'fantastic', 'positive', 'love', 'like', 'happy']
            negative_words = ['bad', 'terrible', 'awful', 'horrible', 'negative', 'hate', 'dislike', 'sad', 'angry', 'disappointed']
            
            text_lower = text.lower()
            positive_count = sum(1 for word in positive_words if word in text_lower)
            negative_count = sum(1 for word in negative_words if word in text_lower)
            
            total_count = positive_count + negative_count
            if total_count == 0:
                return 0.0, "neutral"
            
            sentiment_score = (positive_count - negative_count) / total_count
            
            if sentiment_score > 0.1:
                sentiment_label = "positive"
            elif sentiment_score < -0.1:
                sentiment_label = "negative"
            else:
                sentiment_label = "neutral"
            
            return sentiment_score, sentiment_label
            
        except Exception as e:
            self.logger.error(f"Sentiment analysis failed: {e}")
            return 0.0, "neutral"
    
    async def _classify_content(self, text: str) -> str:
        """Classify content type"""
        text_lower = text.lower()
        
        # Simple classification based on keywords
        if any(word in text_lower for word in ['contract', 'agreement', 'terms', 'conditions', 'legal']):
            return "legal"
        elif any(word in text_lower for word in ['research', 'study', 'analysis', 'methodology', 'conclusion']):
            return "academic"
        elif any(word in text_lower for word in ['manual', 'instructions', 'guide', 'tutorial', 'how to']):
            return "technical"
        elif any(word in text_lower for word in ['report', 'summary', 'findings', 'results', 'executive']):
            return "business"
        elif any(word in text_lower for word in ['story', 'chapter', 'character', 'plot', 'narrative']):
            return "creative"
        else:
            return "general"
    
    async def _generate_summary(self, text: str) -> Optional[str]:
        """Generate summary of text content"""
        try:
            # Simple extractive summarization
            if not TEXT_ANALYSIS_AVAILABLE:
                return None
            
            sentences = sent_tokenize(text)
            if len(sentences) <= 3:
                return text
            
            # Get first few sentences as summary
            summary_sentences = sentences[:min(3, len(sentences))]
            summary = ' '.join(summary_sentences)
            
            # Limit length
            if len(summary) > self.config.max_summary_length:
                summary = summary[:self.config.max_summary_length] + "..."
            
            return summary
            
        except Exception as e:
            self.logger.error(f"Summary generation failed: {e}")
            return None
    
    async def _calculate_quality_metrics(self, text: str) -> Dict[str, float]:
        """Calculate text quality metrics"""
        try:
            metrics = {}
            
            # Basic metrics
            word_count = len(text.split())
            char_count = len(text)
            
            # Text density
            if char_count > 0:
                metrics['text_density'] = word_count / char_count
            else:
                metrics['text_density'] = 0.0
            
            # Average word length
            words = text.split()
            if words:
                avg_word_length = sum(len(word) for word in words) / len(words)
                metrics['avg_word_length'] = avg_word_length
            else:
                metrics['avg_word_length'] = 0.0
            
            # Sentence count and average length
            if TEXT_ANALYSIS_AVAILABLE:
                sentences = sent_tokenize(text)
                metrics['sentence_count'] = len(sentences)
                if sentences:
                    metrics['avg_sentence_length'] = word_count / len(sentences)
                else:
                    metrics['avg_sentence_length'] = 0.0
            
            # Vocabulary richness
            unique_words = set(word.lower() for word in words if word.isalpha())
            if words:
                metrics['vocabulary_richness'] = len(unique_words) / len([w for w in words if w.isalpha()])
            else:
                metrics['vocabulary_richness'] = 0.0
            
            return metrics
            
        except Exception as e:
            self.logger.error(f"Quality metrics calculation failed: {e}")
            return {}
    
    async def _perform_ocr(self, data: bytes, metadata: DocumentMetadata) -> Optional[str]:
        """Perform OCR on image-based documents"""
        try:
            if not self._ocr_available or metadata.format != DocumentFormat.PDF:
                return None
            
            # Convert PDF pages to images and perform OCR
            if ADVANCED_DOC_LIBS_AVAILABLE:
                doc = fitz.open(stream=data, filetype="pdf")
                ocr_text = []
                
                for page_num in range(min(len(doc), 10)):  # Limit OCR to first 10 pages
                    page = doc[page_num]
                    
                    # Convert page to image
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    
                    # Perform OCR
                    image = Image.open(io.BytesIO(img_data))
                    text = pytesseract.image_to_string(image, lang=self.config.ocr_language)
                    
                    if text.strip():
                        ocr_text.append(text)
                
                doc.close()
                return '\n\n'.join(ocr_text)
            
            return None
            
        except Exception as e:
            self.logger.error(f"OCR processing failed: {e}")
            return None
    
    async def _extract_structured_data(
        self,
        data: bytes,
        metadata: DocumentMetadata,
        features: DocumentFeatures
    ) -> Optional[Dict[str, Any]]:
        """Extract structured data from document"""
        try:
            structured_data = {
                "document_info": {
                    "format": metadata.format.value if metadata.format else None,
                    "type": metadata.document_type.value if metadata.document_type else None,
                    "size": metadata.file_size,
                    "pages": metadata.page_count
                },
                "content_summary": {
                    "word_count": metadata.word_count,
                    "character_count": metadata.character_count,
                    "language": metadata.language
                }
            }
            
            # Add structure information
            if features.document_structure:
                structured_data["structure"] = {
                    "sections": len(features.document_structure.sections),
                    "headings": len(features.document_structure.headings),
                    "paragraphs": len(features.document_structure.paragraphs),
                    "tables": len(features.document_structure.tables),
                    "images": len(features.document_structure.images)
                }
            
            # Add content analysis
            if features.content_analysis:
                structured_data["analysis"] = {
                    "topics": features.content_analysis.main_topics,
                    "sentiment": features.content_analysis.sentiment_label,
                    "complexity": features.content_analysis.complexity_level,
                    "classification": features.content_analysis.content_classification
                }
            
            return structured_data
            
        except Exception as e:
            self.logger.error(f"Structured data extraction failed: {e}")
            return None
    
    async def _generate_fingerprint(self, data: bytes) -> str:
        """Generate document fingerprint"""
        try:
            # Create hash of document content
            fingerprint = hashlib.sha256(data).hexdigest()[:32]
            return fingerprint
        except Exception as e:
            self.logger.error(f"Fingerprint generation failed: {e}")
            return ""
    
    async def _generate_content_hash(self, content: str) -> str:
        """Generate hash of extracted content"""
        try:
            # Normalize content and create hash
            normalized_content = ' '.join(content.split()).lower()
            content_hash = hashlib.md5(normalized_content.encode('utf-8')).hexdigest()[:16]
            return content_hash
        except Exception as e:
            self.logger.error(f"Content hash generation failed: {e}")
            return ""
    
    async def _generate_tags(
        self,
        metadata: DocumentMetadata,
        features: DocumentFeatures
    ) -> List[str]:
        """Generate relevant tags for the document"""
        tags = []
        
        try:
            # Format tags
            if metadata.format:
                tags.append(f"format-{metadata.format.value}")
            
            if metadata.document_type:
                tags.append(f"type-{metadata.document_type.value}")
            
            # Size tags
            if metadata.file_size:
                if metadata.file_size < 1024 * 1024:  # < 1MB
                    tags.append("size-small")
                elif metadata.file_size < 10 * 1024 * 1024:  # < 10MB
                    tags.append("size-medium")
                else:
                    tags.append("size-large")
            
            # Language tags
            if metadata.language:
                tags.append(f"lang-{metadata.language}")
            
            # Content tags
            if features.content_analysis:
                if features.content_analysis.content_classification:
                    tags.append(f"content-{features.content_analysis.content_classification}")
                
                if features.content_analysis.complexity_level:
                    tags.append(f"complexity-{features.content_analysis.complexity_level}")
                
                if features.content_analysis.sentiment_label:
                    tags.append(f"sentiment-{features.content_analysis.sentiment_label}")
                
                # Topic tags
                for topic in features.content_analysis.main_topics[:3]:  # Limit to 3 topics
                    tags.append(f"topic-{topic}")
            
            # Structure tags
            if features.document_structure:
                if features.document_structure.tables:
                    tags.append("has-tables")
                if features.document_structure.images:
                    tags.append("has-images")
                if features.document_structure.headings:
                    tags.append("has-headings")
            
            # Processing tags
            if features.ocr_text:
                tags.append("ocr-processed")
            
            return tags
            
        except Exception as e:
            self.logger.error(f"Tag generation failed: {e}")
            return []
    
    async def health_check(self) -> Dict[str, Any]:
        """Perform health check on the document processor"""
        return {
            "status": "healthy" if self._initialized else "not_initialized",
            "doc_libs_available": DOC_LIBS_AVAILABLE,
            "advanced_doc_libs_available": ADVANCED_DOC_LIBS_AVAILABLE,
            "ocr_available": self._ocr_available,
            "text_analysis_available": TEXT_ANALYSIS_AVAILABLE,
            "supported_formats": [fmt.value for fmt in self._supported_formats],
            "config": self.config.__dict__
        }


async def create_document_processor(
    db_session,
    redis_client,
    config: Optional[Dict[str, Any]] = None
) -> DocumentProcessor:
    """
    Factory function to create and initialize a document processor
    
    Args:
        db_session: Database session
        redis_client: Redis client
        config: Configuration dictionary
        
    Returns:
        Initialized DocumentProcessor instance
    """
    # Create config from dict if provided
    processor_config = None
    if config:
        processor_config = DocumentProcessingConfig(**{
            k: v for k, v in config.items() 
            if k in DocumentProcessingConfig.__dataclass_fields__
        })
    
    # Create processor
    processor = DocumentProcessor(
        db_session=db_session,
        redis_client=redis_client,
        config=processor_config
    )
    
    # Initialize
    await processor.initialize()
    
    return processor
